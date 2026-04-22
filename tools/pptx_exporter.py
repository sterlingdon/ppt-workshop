import json
import math
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def is_transparent_color(color_str):
    """Returns True if color is transparent/none (should not be rendered as a fill)."""
    if not color_str:
        return True
    s = color_str.strip().lower()
    if s in ("transparent", "rgba(0, 0, 0, 0)", "rgba(0,0,0,0)"):
        return True
    if s.startswith("rgba("):
        try:
            parts = s[5:].rstrip(")").split(",")
            if len(parts) == 4 and float(parts[3].strip()) < 0.05:
                return True
        except Exception:
            pass
    return False

def parse_px(px_str):
    if not px_str: return 0
    return float(px_str.replace("px", "").strip())

def _linear_to_srgb(c):
    c = max(0.0, min(1.0, c))
    if c <= 0.0031308:
        return c * 12.92
    return 1.055 * (c ** (1.0 / 2.4)) - 0.055

def oklab_to_rgb(L, a, b):
    """oklab(L a b) → (r, g, b) each 0-255."""
    l_ = L + 0.3963377774 * a + 0.2158037573 * b
    m_ = L - 0.1055613458 * a - 0.0638541728 * b
    s_ = L - 0.0894841775 * a - 1.2914855480 * b

    l = l_ ** 3
    m = m_ ** 3
    s = s_ ** 3

    r_lin =  4.0767416621 * l - 3.3077115913 * m + 0.2309699292 * s
    g_lin = -1.2684380046 * l + 2.6097574011 * m - 0.3413193965 * s
    b_lin = -0.0041960863 * l - 0.7034186147 * m + 1.7076147010 * s

    return (
        int(round(_linear_to_srgb(r_lin) * 255)),
        int(round(_linear_to_srgb(g_lin) * 255)),
        int(round(_linear_to_srgb(b_lin) * 255)),
    )

def oklch_to_rgb(L, C, H):
    """oklch(L C H) → (r, g, b) each 0-255."""
    h_rad = math.radians(H)
    return oklab_to_rgb(L, C * math.cos(h_rad), C * math.sin(h_rad))

def parse_rgb(color_str):
    """Parse any CSS color string → (r, g, b).  Defaults to white on failure."""
    if not color_str:
        return 255, 255, 255
    color_str = color_str.strip()

    # oklch(L C H [/ alpha])
    if color_str.startswith("oklch("):
        try:
            inner = color_str[6:].rstrip(")")
            if "/" in inner:
                inner, alpha_s = inner.split("/", 1)
                if float(alpha_s.strip()) < 0.05:
                    return 255, 255, 255  # near-transparent gradient text → white
            parts = inner.strip().split()
            return oklch_to_rgb(float(parts[0]), float(parts[1]), float(parts[2]))
        except Exception:
            return 255, 255, 255

    # oklab(L a b [/ alpha])
    if color_str.startswith("oklab("):
        try:
            inner = color_str[6:].rstrip(")")
            if "/" in inner:
                inner, alpha_s = inner.split("/", 1)
                if float(alpha_s.strip()) < 0.05:
                    return 240, 240, 240
            parts = inner.strip().split()
            return oklab_to_rgb(float(parts[0]), float(parts[1]), float(parts[2]))
        except Exception:
            return 255, 255, 255

    # rgb(...) / rgba(...)
    try:
        s = color_str.replace("rgb(", "").replace("rgba(", "").replace(")", "")
        parts = [float(p.strip()) for p in s.split(",")]
        if len(parts) == 4 and parts[3] < 0.05:
            return 255, 255, 255  # near-transparent = bg-clip-text gradient → pure white
        return int(parts[0]), int(parts[1]), int(parts[2])
    except Exception:
        return 255, 255, 255  # unknown format → white (safe on dark backgrounds)

# Maps CSS system-font aliases (unrecognised by PPT) to usable fallbacks
_SYSTEM_FONT_MAP = {
    "ui-sans-serif":     "Arial",
    "system-ui":         "Arial",
    "-apple-system":     "Arial",
    "BlinkMacSystemFont":"Arial",
    "ui-serif":          "Georgia",
    "ui-monospace":      "Consolas",
    "ui-rounded":        "Arial",
}

def clean_font_name(font_string):
    """'Inter', sans-serif  →  'Inter'.  System aliases  →  known PPT font."""
    try:
        first = font_string.split(",")[0].strip().strip("'").strip('"')
        return _SYSTEM_FONT_MAP.get(first, first) or "Arial"
    except Exception:
        return "Arial"

def build_pptx(manifest_path, output_path):
    print(f"[*] 读取构件化坐标清单: {manifest_path}")
    with open(manifest_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width  = Emu(12192000)   # 13.333 in  (16:9)
    prs.slide_height = Emu(6858000)    # 7.5   in
    blank_layout = prs.slide_layouts[6]

    # CSS-pixel → EMU scale (viewport 1920 px maps to full slide width)
    SCALE = prs.slide_width / 1920     # ≈ 6350 EMU/px

    def _apply_run_style(run, run_style, fallback_style):
        fw_str = str(run_style.get("fontWeight") or fallback_style.get("fontWeight", "400")).strip()
        try:
            fw_num = int(fw_str)
        except ValueError:
            fw_num = 700 if fw_str == "bold" else 400

        run.font.name  = clean_font_name(run_style.get("fontFamily") or fallback_style.get("fontFamily", "Arial"))
        pt_size        = parse_px(run_style.get("fontSize") or fallback_style.get("fontSize", "16px")) * SCALE / 12700
        run.font.size  = Pt(int(round(pt_size)))
        r, g, b        = parse_rgb(run_style.get("color") or fallback_style.get("color", "rgb(255,255,255)"))
        run.font.color.rgb = RGBColor(r, g, b)
        if fw_num >= 600:
            run.font.bold = True
        if (run_style.get("fontStyle") or fallback_style.get("fontStyle", "")) == "italic":
            run.font.italic = True
        deco = (run_style.get("textDecorationLine") or fallback_style.get("textDecorationLine", "")).lower()
        if "underline" in deco:
            run.font.underline = True
        if "line-through" in deco:
            rPr = run._r.get_or_add_rPr()
            rPr.set("strike", "sngStrike")

    def _pptx_align(css_align):
        if css_align == "center": return PP_ALIGN.CENTER
        if css_align in ("right", "end"): return PP_ALIGN.RIGHT
        return PP_ALIGN.LEFT

    def _add_text(slide, txt, offset_x=0, offset_y=0):
        box = txt["box"]
        style = txt["style"]
        text_content = style.get("text", "").strip()
        if not text_content:
            return

        x_emu = int((box["x"] + offset_x) * SCALE)
        y_emu = int((box["y"] + offset_y) * SCALE)
        h_emu = int(box["height"] * SCALE)

        raw_font_px = parse_px(style.get("fontSize", "16px"))
        exact_w_emu = int(box["width"] * SCALE)
        # Subtract CSS padding from height before single-line check: labels have
        # vertical padding that inflates bbox height past the threshold.
        pad_v_px = (parse_px(style.get("paddingTop", "0px"))
                    + parse_px(style.get("paddingBottom", "0px")))
        content_h_emu = h_emu - int(pad_v_px * SCALE)
        is_singleline = content_h_emu <= int(raw_font_px * SCALE * 1.8)
        align_css = style.get("textAlign", "left")

        # Background fill shape for labels/tags.
        bg_color_str = style.get("backgroundColor", "")
        has_bg = not is_transparent_color(bg_color_str)
        if has_bg:
            bg_r, bg_g, bg_b = parse_rgb(bg_color_str)
            br_str = style.get("borderRadius", "0px")
            has_radius = br_str and br_str not in ("0px", "0%", "")
            shape_type = (MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
                          if has_radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE)
            bg_shape = slide.shapes.add_shape(shape_type, x_emu, y_emu, exact_w_emu, h_emu)
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)
            bg_shape.line.fill.background()

        # Labels: never expand width. Plain text: expand non-center single-line
        # boxes to absorb font-metric differences.
        if has_bg or not is_singleline:
            w_emu = exact_w_emu
            word_wrap = not is_singleline
        else:
            if align_css == "center":
                w_emu = exact_w_emu
            else:
                max_w = int(prs.slide_width) - x_emu
                w_emu = min(int(exact_w_emu * 1.4), max_w)
            word_wrap = False

        tx_box = slide.shapes.add_textbox(x_emu, y_emu, w_emu, h_emu)
        tf = tx_box.text_frame
        tf.margin_top = int(parse_px(style.get("paddingTop", "0px")) * SCALE)
        tf.margin_bottom = int(parse_px(style.get("paddingBottom", "0px")) * SCALE)
        tf.margin_left = int(parse_px(style.get("paddingLeft", "0px")) * SCALE)
        tf.margin_right = int(parse_px(style.get("paddingRight", "0px")) * SCALE)
        tf.word_wrap = word_wrap
        if has_bg:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        runs_data = style.get("runs")
        if runs_data:
            paragraphs = []
            cur = []
            for run_item in runs_data:
                if run_item["text"] == "\n":
                    paragraphs.append(cur)
                    cur = []
                elif run_item["text"]:
                    cur.append(run_item)
            if cur:
                paragraphs.append(cur)

            for para_idx, para_runs in enumerate(paragraphs):
                p = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
                p.alignment = _pptx_align(align_css)
                for run_item in para_runs:
                    run = p.add_run()
                    run.text = run_item["text"]
                    _apply_run_style(run, run_item.get("style", {}), style)
        else:
            tf.text = text_content
            p = tf.paragraphs[0]
            p.alignment = _pptx_align(align_css)
            run = p.runs[0]
            _apply_run_style(run, style, style)

    for slide_data in data.get("slides", []):
        idx = slide_data["index"] + 1
        print(f"[*] 正在编排幻灯片第 {idx} 页...")
        slide = prs.slides.add_slide(blank_layout)

        # ── 1. 全页底图 ──────────────────────────────────────────────
        bg_path = slide_data.get("bg_path")
        if bg_path:
            slide.shapes.add_picture(bg_path, 0, 0, prs.slide_width, prs.slide_height)

        # ── 2. 局部组件截图（GlassCard 等） ──────────────────────────
        for comp in slide_data.get("components", []):
            box = comp["box"]
            slide.shapes.add_picture(
                comp["path"],
                int(box["x"]      * SCALE),
                int(box["y"]      * SCALE),
                int(box["width"]  * SCALE),
                int(box["height"] * SCALE),
            )

        # ── 3. Item-aware groups ────────────────────────────────────────
        for group in slide_data.get("groups", []):
            for segment in group.get("segments", []):
                raster = segment.get("raster", {})
                path = raster.get("path")
                box = segment.get("box")
                if path and box:
                    slide.shapes.add_picture(
                        path,
                        int(box["x"] * SCALE),
                        int(box["y"] * SCALE),
                        int(box["width"] * SCALE),
                        int(box["height"] * SCALE),
                    )

            for item in group.get("items", []):
                box = item["box"]
                raster = item.get("raster", {})
                path = raster.get("path")
                if path:
                    slide.shapes.add_picture(
                        path,
                        int(box["x"] * SCALE),
                        int(box["y"] * SCALE),
                        int(box["width"] * SCALE),
                        int(box["height"] * SCALE),
                    )
                for txt in item.get("texts", []):
                    _add_text(slide, txt, offset_x=box["x"], offset_y=box["y"])

        # ── 4. 原生可编辑文字层 ──────────────────────────────────────
        for txt in slide_data.get("texts", []):
            _add_text(slide, txt)

    prs.save(output_path)
    print(f"[*] 转化完成！保存至: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python pptx_exporter.py <manifest.json> <output.pptx>")
        sys.exit(1)
    build_pptx(sys.argv[1], sys.argv[2])
