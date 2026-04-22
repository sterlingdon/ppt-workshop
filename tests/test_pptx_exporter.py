import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))
from tools.pptx_exporter import build_pptx


def write_manifest(tmp_path, slides):
    manifest = tmp_path / "layout_manifest.json"
    manifest.write_text(json.dumps({"slides": slides}), encoding="utf-8")
    return manifest


def test_manifest_export_creates_correct_slide_count(tmp_path):
    manifest = write_manifest(
        tmp_path,
        [
            {"index": 0, "texts": [], "components": [], "bg_path": None},
            {"index": 1, "texts": [], "components": [], "bg_path": None},
            {"index": 2, "texts": [], "components": [], "bg_path": None},
        ],
    )
    output = str(tmp_path / "test.pptx")
    build_pptx(str(manifest), output)
    prs = Presentation(output)
    assert len(prs.slides) == 3


def test_manifest_export_slides_are_widescreen(tmp_path):
    manifest = write_manifest(tmp_path, [{"index": 0, "texts": [], "components": [], "bg_path": None}])
    output = str(tmp_path / "test.pptx")
    build_pptx(str(manifest), output)
    prs = Presentation(output)
    assert abs(prs.slide_width - Inches(13.333)) < Inches(0.01)
    assert abs(prs.slide_height - Inches(7.5)) < Inches(0.01)


def test_manifest_export_adds_native_text(tmp_path):
    manifest = write_manifest(
        tmp_path,
        [
            {
                "index": 0,
                "texts": [
                    {
                        "box": {"x": 120, "y": 160, "width": 500, "height": 80},
                        "style": {
                            "text": "AI的未来",
                            "fontSize": "48px",
                            "fontFamily": "Inter, sans-serif",
                            "color": "rgb(255,255,255)",
                            "fontWeight": "800",
                            "textAlign": "left",
                        },
                    }
                ],
                "components": [],
                "bg_path": None,
            }
        ],
    )
    output = str(tmp_path / "test.pptx")
    build_pptx(str(manifest), output)
    prs = Presentation(output)
    slide = prs.slides[0]
    all_text = " ".join(shape.text for shape in slide.shapes if shape.has_text_frame)
    assert "AI的未来" in all_text


def test_manifest_export_accepts_bg_images(tmp_path):
    from PIL import Image

    bg_img = tmp_path / "bg_0.png"
    img = Image.new("RGB", (1920, 1080), color=(10, 8, 24))
    img.save(str(bg_img))
    manifest = write_manifest(
        tmp_path,
        [{"index": 0, "texts": [], "components": [], "bg_path": str(bg_img)}],
    )
    output = str(tmp_path / "with_bg.pptx")
    build_pptx(str(manifest), output)
    prs = Presentation(output)
    slide = prs.slides[0]
    assert len(slide.shapes) > 0


def test_manifest_export_adds_item_raster_and_item_text(tmp_path):
    from PIL import Image

    item_img = tmp_path / "item.png"
    Image.new("RGBA", (400, 120), color=(20, 20, 20, 255)).save(str(item_img))

    manifest = write_manifest(
        tmp_path,
        [
            {
                "index": 0,
                "texts": [],
                "components": [],
                "groups": [
                    {
                        "id": "group_0_0",
                        "kind": "list",
                        "mode": "item-hybrid",
                        "box": {"x": 100, "y": 100, "width": 500, "height": 200},
                        "items": [
                            {
                                "id": "item_0_0_0",
                                "mode": "item-hybrid",
                                "box": {"x": 120, "y": 140, "width": 400, "height": 120},
                                "raster": {"path": str(item_img), "mode": "item"},
                                "texts": [
                                    {
                                        "box": {"x": 80, "y": 30, "width": 250, "height": 40},
                                        "style": {
                                            "text": "独立列表项",
                                            "fontSize": "24px",
                                            "fontFamily": "Inter, sans-serif",
                                            "color": "rgb(255,255,255)",
                                            "fontWeight": "700",
                                            "textAlign": "left",
                                        },
                                    }
                                ],
                                "bullets": [],
                            }
                        ],
                        "segments": [],
                    }
                ],
                "bg_path": None,
            }
        ],
    )

    output = str(tmp_path / "item_group.pptx")
    build_pptx(str(manifest), output)
    prs = Presentation(output)
    slide = prs.slides[0]
    all_text = " ".join(shape.text for shape in slide.shapes if shape.has_text_frame)
    assert "独立列表项" in all_text
    assert len(slide.shapes) >= 2
